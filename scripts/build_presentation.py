#!/usr/bin/env python3
"""Build the Huatuo Big Data Processing presentation deck."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "presentations" / "huatuo_big_data_processing_20min.pptx"

WIDE = (13.333, 7.5)

COLORS = {
    "bg": RGBColor(247, 249, 250),
    "ink": RGBColor(24, 35, 45),
    "muted": RGBColor(91, 105, 116),
    "blue": RGBColor(47, 111, 159),
    "orange": RGBColor(217, 119, 6),
    "green": RGBColor(38, 126, 101),
    "line": RGBColor(214, 222, 229),
    "pale_blue": RGBColor(230, 241, 249),
    "pale_orange": RGBColor(252, 238, 221),
    "white": RGBColor(255, 255, 255),
}


def set_slide_bg(slide, color=COLORS["bg"]) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int = 24,
    color=COLORS["ink"],
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    font: str = "Aptos",
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0)
    frame.margin_right = Inches(0)
    frame.margin_top = Inches(0)
    frame.margin_bottom = Inches(0)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_text(slide, title, 0.72, 0.42, 8.8, 0.58, size=27, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.74, 1.02, 8.8, 0.34, size=13, color=COLORS["muted"])


def add_footer(slide, number: int) -> None:
    add_text(
        slide,
        "Huatuo Big Data Processing",
        0.72,
        7.05,
        3.2,
        0.2,
        size=8,
        color=COLORS["muted"],
    )
    add_text(
        slide,
        f"{number:02d}",
        12.08,
        7.03,
        0.45,
        0.2,
        size=8,
        color=COLORS["muted"],
        align=PP_ALIGN.RIGHT,
    )


def add_rule(slide, x: float, y: float, w: float, color=COLORS["line"]) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.015)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_bullets(
    slide,
    bullets: list[str],
    x: float,
    y: float,
    w: float,
    h: float,
    size: int = 19,
    color=COLORS["ink"],
    gap: float = 0.2,
) -> None:
    top = y
    for bullet in bullets:
        add_text(slide, "-", x, top + 0.01, 0.18, 0.22, size=size, color=COLORS["blue"], bold=True)
        add_text(slide, bullet, x + 0.28, top, w - 0.28, 0.35, size=size, color=color)
        top += gap + 0.35


def add_metric(slide, label: str, value: str, x: float, y: float, w: float, accent) -> None:
    add_text(slide, value, x, y, w, 0.48, size=30, color=accent, bold=True)
    add_text(slide, label, x, y + 0.56, w, 0.28, size=12, color=COLORS["muted"])


def add_pill(slide, text: str, x: float, y: float, w: float, color) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.38)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    add_text(slide, text, x + 0.12, y + 0.095, w - 0.24, 0.16, size=10, color=COLORS["ink"], bold=True, align=PP_ALIGN.CENTER)


def add_table(slide, rows: list[list[str]], x: float, y: float, w: float, h: float, widths: list[float] | None = None):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if widths:
        for index, width in enumerate(widths):
            table.columns[index].width = Inches(width)
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = COLORS["pale_blue"] if r == 0 else COLORS["white"]
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.name = "Aptos"
                    run.font.size = Pt(10 if len(rows) > 5 else 12)
                    run.font.bold = r == 0
                    run.font.color.rgb = COLORS["ink"]
    return table_shape


def add_native_benchmark_chart(slide, x: float, y: float, w: float, h: float) -> None:
    data = CategoryChartData()
    data.categories = [
        "Dept demand",
        "Disease trends",
        "QA quality",
        "Question type",
    ]
    data.add_series("Hadoop", (128.071, 139.457, 273.549, 204.027))
    data.add_series("Spark", (117.327, 149.503, 82.296, 113.497))
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
        data,
    )
    chart = graphic_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.TOP
    chart.legend.include_in_layout = False
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.axis_title.text_frame.text = "Seconds"
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.font.size = Pt(9)
    for series, color in zip(chart.series, (COLORS["blue"], COLORS["orange"])):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = color


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(WIDE[0])
    prs.slide_height = Inches(WIDE[1])
    blank = prs.slide_layouts[6]

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, COLORS["ink"])
    add_text(slide, "34M", 0.82, 0.82, 3.1, 0.82, size=48, color=COLORS["orange"], bold=True)
    add_text(slide, "medical QA records", 0.88, 1.66, 3.4, 0.28, size=15, color=COLORS["white"])
    add_text(
        slide,
        "Huatuo Big Data Processing",
        0.82,
        3.05,
        9.6,
        0.62,
        size=34,
        color=COLORS["white"],
        bold=True,
    )
    add_text(
        slide,
        "Schema normalization, Hadoop Streaming, and Apache Spark analytics on HDFS/YARN",
        0.86,
        3.82,
        8.8,
        0.45,
        size=18,
        color=RGBColor(213, 224, 232),
    )
    add_text(slide, "WOC7017 Big Data Processing | 20-minute presentation", 0.88, 6.62, 6.6, 0.24, size=12, color=RGBColor(185, 199, 210))
    for x, y, w, color in [(9.2, 0.7, 2.7, COLORS["blue"]), (10.45, 2.0, 1.9, COLORS["green"]), (8.8, 4.95, 3.1, COLORS["orange"])]:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.2))
        shape.rotation = -18
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

    # Slide content specs
    specs = [
        (
            "20-minute route",
            "A focused walkthrough from data problem to measured framework behavior.",
            [
                ("2 min", "Problem and dataset scale"),
                ("4 min", "Integration method and canonical schema"),
                ("4 min", "Pipeline and MapReduce case studies"),
                ("4 min", "Pilot and full-data validation"),
                ("4 min", "Hadoop/Spark benchmark results"),
                ("2 min", "Takeaways and limitations"),
            ],
        ),
        (
            "Problem: heterogeneous medical QA at scale",
            "The task is not only analytics; first the data has to become analyzable.",
            [
                "Four Huatuo sources publish related medical QA records with different schemas.",
                "The raw snapshot is 5.432 GB and more than 34 million JSONL records.",
                "The pipeline must preserve provenance and avoid false joins across sources.",
                "The final workload compares Hadoop Streaming and Spark on equivalent tasks.",
            ],
        ),
        (
            "Dataset: four sources, one analytical target",
            "The source mix is highly imbalanced, which affects both storage and analysis.",
            [
                ["Source", "Records", "Size"],
                ["Consultation QA", "32,708,346", "4.537 GB"],
                ["Encyclopedia QA", "364,420", "0.608 GB"],
                ["Knowledge Graph QA", "798,444", "0.149 GB"],
                ["Huatuo26M-Lite", "177,703", "0.138 GB"],
                ["Combined", "34,048,913", "5.432 GB"],
            ],
        ),
        (
            "Integration decision: union, not join",
            "The datasets share a medical QA structure, but not a trustworthy shared entity key.",
            [
                "No universal patient, consultation, or question identifier exists across sources.",
                "Joining would invent relationships that the data does not prove.",
                "Each record is mapped into one canonical schema and kept as a separate observation.",
                "Source provenance remains explicit through source and source_split fields.",
            ],
        ),
        (
            "Canonical schema makes the sources comparable",
            "Normalization is formatting cleanup plus transparent metadata, not medical rewriting.",
            [
                "`question`, `answer`, `answer_type`",
                "`source`, `source_split`, `metadata_origin`",
                "`department`, `related_disease`, `quality_score`",
                "`question_length`, `answer_length`",
                "`question_hash`, `duplicate_flag`",
            ],
        ),
        (
            "Processing pipeline",
            "The same canonical JSONL file becomes the input for both frameworks.",
            [
                "Download original Hugging Face snapshots",
                "Create source integrity manifest",
                "Build pilot subset for validation",
                "Standardize and union all source records",
                "Inspect full canonical data",
                "Upload to HDFS and run Hadoop/Spark jobs on YARN",
            ],
        ),
        (
            "Four MapReduce-style case studies",
            "Each analytical question is implemented with equivalent Hadoop and Spark logic.",
            [
                ["Case", "Map behavior", "Reduce / aggregate behavior"],
                ["1. Department demand", "Emit observed department", "Count by department"],
                ["2. Disease/symptom trends", "Emit matched terms", "Count by term"],
                ["3. QA quality", "Emit source quality record", "Aggregate averages and rates"],
                ["4. Question type", "Emit classified category", "Count by category"],
            ],
        ),
        (
            "Pilot validation before full-scale execution",
            "A deterministic 400,000-record pilot reduced risk before running 34M records.",
            [
                "100,000 real records sampled from each source.",
                "Pipeline validation confirmed 400,000 canonical output records.",
                "Hadoop and Spark produced matching output group counts on all case studies.",
                "Pilot timing was useful for development, not the final performance conclusion.",
            ],
        ),
        (
            "Full canonical dataset inspection",
            "The final unified file is larger than the raw snapshot because each row carries metadata.",
            [
                ("34,048,898", "canonical records"),
                ("19 GB", "final JSONL size"),
                ("29.847%", "duplicate-flagged rows"),
                ("0", "invalid JSON lines"),
            ],
        ),
        (
            "Benchmark design",
            "The comparison uses the same HDFS input and one timed job per case study.",
            [
                "Input: /healthcare/full/input/huatuo_unified.jsonl",
                "Hadoop: Hadoop Streaming mapper/reducer submitted to YARN.",
                "Spark: PySpark application submitted to YARN.",
                "Timing: elapsed wall-clock seconds per framework per case study.",
                "Validation: output row counts logged and compared.",
            ],
        ),
        (
            "Full benchmark result",
            "Spark wins 3 of 4 full-data case studies; Hadoop remains competitive for streaming text matching.",
            [],
        ),
        (
            "What the benchmark means",
            "Small pilot timing favored Hadoop because Spark startup overhead was large relative to the workload.",
            [
                "On the full 34M-record dataset, Spark's execution engine becomes more competitive.",
                "Spark is much faster for QA quality because distributed aggregation avoids reducer pressure.",
                "Hadoop is slightly faster for disease/symptom trends, a simple streaming text scan.",
                "The result supports comparing frameworks on realistic workload size, not pilot size alone.",
            ],
        ),
        (
            "Reproducibility and evidence",
            "The project is script-driven so the pipeline can be inspected, rerun, and defended.",
            [
                "Source manifest: reports/source_manifest.json",
                "Full-data summary: reports/full_summary.json",
                "Benchmark CSV: results/case-benchmark/timings-full-cluster-20260607-171724.csv",
                "Chart: reports/full_benchmark_chart.png",
                "Detailed guide: journey.md",
            ],
        ),
    ]

    for index, spec in enumerate(specs, start=2):
        slide = prs.slides.add_slide(blank)
        set_slide_bg(slide)
        add_title(slide, spec[0], spec[1])
        add_rule(slide, 0.72, 1.48, 11.9)

        if index == 2:
            for row, (duration, label) in enumerate(spec[2]):
                y = 2.0 + row * 0.62
                add_text(slide, duration, 0.9, y, 0.9, 0.3, size=18, color=COLORS["orange"], bold=True)
                add_text(slide, label, 1.95, y, 7.6, 0.3, size=18)
            add_text(slide, "Target pacing: about 1-2 minutes per content slide.", 0.9, 6.2, 7.4, 0.28, size=14, color=COLORS["muted"])
        elif index == 4:
            add_table(slide, spec[2], 0.82, 1.9, 11.4, 3.55, widths=[4.1, 2.3, 1.8])
            add_metric(slide, "original source snapshot", "5.432 GB", 0.94, 5.95, 2.4, COLORS["blue"])
            add_metric(slide, "source records", "34.0M", 4.1, 5.95, 2.4, COLORS["orange"])
            add_metric(slide, "dominant source", "96%", 7.2, 5.95, 2.4, COLORS["green"])
        elif index == 6:
            add_bullets(slide, spec[2], 0.9, 1.95, 6.0, 4.8, size=18, gap=0.22)
            code = (
                "{\n"
                '  "source": "consultation",\n'
                '  "question": "...",\n'
                '  "answer_type": "text/url",\n'
                '  "question_hash": "sha256",\n'
                '  "duplicate_flag": false\n'
                "}"
            )
            box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.25), Inches(1.88), Inches(4.7), Inches(3.55))
            box.fill.solid()
            box.fill.fore_color.rgb = COLORS["ink"]
            box.line.fill.background()
            add_text(slide, code, 7.55, 2.16, 4.1, 2.65, size=16, color=COLORS["white"], font="Menlo")
            add_text(slide, "Original meaning is preserved; only formatting and metadata are standardized.", 7.3, 5.72, 4.7, 0.46, size=14, color=COLORS["muted"])
        elif index == 7:
            steps = spec[2]
            x_positions = [0.85, 2.85, 4.85, 6.85, 8.85, 10.85]
            for step_num, (x, label) in enumerate(zip(x_positions, steps), start=1):
                add_text(slide, f"{step_num}", x, 2.0, 0.38, 0.3, size=17, color=COLORS["orange"], bold=True, align=PP_ALIGN.CENTER)
                shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x - 0.05), Inches(2.42), Inches(1.25), Inches(0.06))
                shape.fill.solid()
                shape.fill.fore_color.rgb = COLORS["blue"] if step_num < 6 else COLORS["green"]
                shape.line.fill.background()
                add_text(slide, label, x - 0.25, 2.75, 1.65, 1.05, size=12, color=COLORS["ink"], align=PP_ALIGN.CENTER)
            add_text(slide, "Source JSONL -> canonical JSONL -> HDFS input -> framework jobs -> result evidence", 1.35, 5.55, 10.2, 0.36, size=18, color=COLORS["muted"], align=PP_ALIGN.CENTER)
        elif index == 8:
            add_table(slide, spec[2], 0.72, 1.9, 11.85, 3.95, widths=[2.5, 4.25, 4.55])
        elif index == 10:
            for col, (value, label) in enumerate(spec[2]):
                x = 1.05 + col * 2.95
                add_metric(slide, label, value, x, 2.15, 2.45, [COLORS["blue"], COLORS["orange"], COLORS["green"], COLORS["blue"]][col])
            add_text(slide, "Inspection created reports/full_summary.json and source-specific canonical samples.", 1.05, 4.55, 8.9, 0.38, size=17)
            add_text(slide, "The processed file is larger than the raw snapshot because every row carries provenance, length, quality, hash, and duplicate metadata.", 1.05, 5.15, 10.2, 0.62, size=16, color=COLORS["muted"])
        elif index == 12:
            add_native_benchmark_chart(slide, 0.85, 1.82, 7.4, 4.35)
            rows = [
                ["Case", "Hadoop", "Spark", "Faster"],
                ["Dept demand", "128.1s", "117.3s", "Spark"],
                ["Disease trends", "139.5s", "149.5s", "Hadoop"],
                ["QA quality", "273.5s", "82.3s", "Spark"],
                ["Question type", "204.0s", "113.5s", "Spark"],
            ]
            add_table(slide, rows, 8.55, 1.96, 3.95, 3.25, widths=[1.35, 0.85, 0.85, 0.9])
            add_text(slide, "Full command: bash scripts/run_full_case_benchmark.sh", 8.62, 5.72, 3.8, 0.32, size=12, color=COLORS["muted"])
        elif index in (3, 5, 9, 11, 13, 14):
            add_bullets(slide, spec[2], 0.9, 1.95, 10.4, 4.8, size=18, gap=0.24)
            if index == 5:
                add_pill(slide, "UNION", 8.4, 5.6, 1.25, COLORS["pale_blue"])
                add_pill(slide, "NOT JOIN", 9.85, 5.6, 1.55, COLORS["pale_orange"])
            if index == 9:
                add_metric(slide, "pilot input", "400K", 8.65, 5.35, 1.8, COLORS["blue"])
                add_metric(slide, "sources", "4", 10.55, 5.35, 1.0, COLORS["orange"])
        else:
            add_bullets(slide, spec[2], 0.9, 1.95, 10.4, 4.8, size=18, gap=0.24)

        add_footer(slide, index - 1)

    return prs


def main() -> int:
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs = build_deck()
    prs.save(OUTPUT)
    print(OUTPUT)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
